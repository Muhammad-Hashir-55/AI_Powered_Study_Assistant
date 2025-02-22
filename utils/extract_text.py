from PyPDF2 import PdfReader
from pptx import Presentation

def extract_text_from_file(filepath):
    if filepath.endswith('.pdf'):
        return extract_text_from_pdf(filepath)
    elif filepath.endswith('.pptx'):
        return extract_text_from_pptx(filepath)
    else:
        raise ValueError("Unsupported file format")

def extract_text_from_pdf(filepath):
    reader = PdfReader(filepath)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text